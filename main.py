from pptx import Presentation
from pptx.util import Inches, Pt #localização no slide e tamanho pra fonte
from pptx.chart.data import CategoryChartData #informações.valores do grafico
from pptx.enum.chart import XL_CHART_TYPE #tipo do grafico

# create presentation
presentation =Presentation()

slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
# slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
# slide3 = presentation.slides.add_slide(presentation.slide_layouts[2])
# slide4 = presentation.slides.add_slide(presentation.slide_layouts[3])
# slide5 = presentation.slides.add_slide(presentation.slide_layouts[4])
# slide6 = presentation.slides.add_slide(presentation.slide_layouts[5])
slide7 = presentation.slides.add_slide(presentation.slide_layouts[6])
# slide8 = presentation.slides.add_slide(presentation.slide_layouts[7])
# slide9 = presentation.slides.add_slide(presentation.slide_layouts[8])

title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Hello World"
subtitle.text = "python-pptx was here!"


# text in slide
# x,y,width,height
textbox = slide7.shapes.add_textbox(Inches(0.5),Inches(0.5),Inches(2),Inches(2))
textbox.text = "Text in a textbox"

text_frame = textbox.text_frame
paragraph = text_frame.add_paragraph()
paragraph.text = "This is some text"
paragraph.font.bold = True
paragraph.font.size = Pt(30)
paragraph = text_frame.add_paragraph()
paragraph.text = "Other text"


# grafico
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

products = ['a','b','c'] #x
prices = [1000,2000,3000] #y
x= Inches(1)
y =Inches(1)
width =Inches(6)
height =Inches(3)


chart_data = CategoryChartData()
chart_data.categories = products
chart_data.add_series(name="Preco", values=prices)

# type of chart,x,y,width,height,dados do grafico
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, width, height, chart_data)

# save presentation
presentation.save('presentation.pptx')